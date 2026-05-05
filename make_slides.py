"""
Generate PriceNest development update presentation as .pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors
NAVY   = RGBColor(0x0B, 0x1C, 0x3E)
YELLOW = RGBColor(0xFD, 0xBB, 0x11)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
SLATE  = RGBColor(0x64, 0x74, 0x8B)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
BLUE   = RGBColor(0x3B, 0x82, 0xF6)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
RED    = RGBColor(0xEF, 0x44, 0x44)
DARKBG = RGBColor(0x0D, 0x23, 0x47)
DARKER = RGBColor(0x06, 0x0E, 0x22)
CARD   = RGBColor(0xF8, 0xFA, 0xFC)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill=NAVY):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def bullet_box(slide, items, x, y, w, h,
               size=14, color=WHITE, dot_color=YELLOW):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(2)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = dot_color
        r1 = p.add_run()
        r1.text = item
        r1.font.size = Pt(size)
        r1.font.color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ─────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, Inches(0.35), H, YELLOW)
add_rect(s, 0, 0, W, Inches(0.08), YELLOW)

add_text(s, "PriceNest",
         Inches(0.7), Inches(1.5), Inches(11), Inches(1.5),
         size=72, bold=True, color=WHITE)
add_text(s, "AI-Powered California Real Estate Valuation Engine",
         Inches(0.7), Inches(3.0), Inches(11), Inches(0.75),
         size=28, bold=False, color=YELLOW)
add_text(s, "April 2026 Development Update",
         Inches(0.7), Inches(3.85), Inches(9), Inches(0.5),
         size=18, color=SLATE, italic=True)
add_text(s, "Jake  ·  Saketh  ·  Aditiya",
         Inches(0.7), Inches(5.1), Inches(9), Inches(0.5),
         size=20, bold=True, color=YELLOW)

add_rect(s, 0, H - Inches(0.5), W, Inches(0.5), DARKER)
add_text(s, "Confidential  ·  April 2026",
         Inches(0.4), H - Inches(0.45), Inches(6), Inches(0.38),
         size=10, color=SLATE, italic=True)

add_notes(s, """• Intro: Jake, Saketh, Aditiya built PriceNest from scratch
• CA-only AI real estate valuation — all 8,000+ census tracts
• This deck covers everything shipped since April 17th""")

# ─────────────────────────────────────────────
# SLIDE 2 — What is PriceNest?
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, LIGHT)
add_rect(s, 0, 0, W, Inches(1.2), NAVY)
add_rect(s, 0, Inches(1.2), Inches(0.07), H - Inches(1.2), YELLOW)

add_text(s, "What is PriceNest?",
         Inches(0.4), Inches(0.18), Inches(10), Inches(0.82),
         size=34, bold=True, color=WHITE)
add_text(s, "Platform Overview",
         Inches(0.4), Inches(0.84), Inches(8), Inches(0.35),
         size=13, color=YELLOW, italic=True)

bullet_box(s, [
    "AI-driven home valuation covering all of California",
    "Instant price estimates with ±7% confidence ranges",
    "5-tab interface: Home, Live Map, Pro Analysis, Tracker, LifeBudget",
    "No login required — runs entirely in the browser",
    "Built on real ACS 2021-22 Census data (~8,000 tracts)",
], Inches(0.5), Inches(1.45), Inches(6.1), Inches(4.8),
   size=15, color=NAVY, dot_color=YELLOW)

bullet_box(s, [
    "Neighborhood Intelligence: income, density, coast & city distances",
    "20-year property value forecast under 3 market scenarios",
    "Mortgage breakdown with PITI, PMI, and break-even vs. rent",
    "Property Tracker with localStorage persistence",
    "Investment score across 5 weighted dimensions",
], Inches(7.0), Inches(1.45), Inches(5.9), Inches(4.8),
   size=15, color=NAVY, dot_color=BLUE)

add_rect(s, 0, H - Inches(0.4), W, Inches(0.4), NAVY)
add_notes(s, """• 5-tab SPA: all features accessible without a page reload
• 8,000+ census tracts cover every corner of California
• No account needed — designed for fast, frictionless analysis""")

# ─────────────────────────────────────────────
# SLIDE 3 — Tech Stack
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.08), YELLOW)

add_text(s, "Tech Stack",
         Inches(0.5), Inches(0.22), Inches(10), Inches(0.82),
         size=36, bold=True, color=WHITE)
add_text(s, "Tools, libraries & services powering PriceNest",
         Inches(0.5), Inches(0.88), Inches(11), Inches(0.4),
         size=14, color=YELLOW, italic=True)

cols = [
    ("Backend", BLUE, [
        "Flask (Python) — REST API + rate limiting",
        "Flask-Limiter: 60 req/min per IP",
        "GeoPandas + Shapely — spatial joins",
        "Scikit-learn — OOF stacking ensemble",
        "XGBoost + LightGBM — gradient boosters",
        "NumPy / Pandas — data pipeline",
        "LRU cache — geocode & API response cache",
    ]),
    ("Data & APIs", YELLOW, [
        "ACS 2021-22 Census (~8,000 CA tracts)",
        "US Census TIGER shapefiles (tract polygons)",
        "Nominatim / OpenStreetMap — geocoding",
        "RentCast AVM — comparable sales values",
        "RentCast Rent AVM — monthly rent estimates",
        "Google Street View Static API — photos",
        "CA coastline KD-tree (EPSG:3310)",
    ]),
    ("Frontend", GREEN, [
        "Tailwind CSS — utility-first styling",
        "Dark mode via html.dark class toggle",
        "Chart.js — buy/rent, forecast, score charts",
        "Leaflet.js — interactive CA heatmap",
        "localStorage — tracker + search history",
        "Intl.NumberFormat — currency formatting",
        "Vanilla JS — no framework overhead",
    ]),
]

for i, (title, accent, items) in enumerate(cols):
    cx = Inches(0.35 + i * 4.3)
    card = s.shapes.add_shape(1, cx, Inches(1.42), Inches(4.1), Inches(5.72))
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = DARKBG
    add_rect(s, cx, Inches(1.42), Inches(4.1), Inches(0.07), accent)
    add_text(s, title, cx + Inches(0.15), Inches(1.54), Inches(3.8), Inches(0.44),
             size=15, bold=True, color=accent)
    bullet_box(s, items, cx + Inches(0.12), Inches(2.02), Inches(3.85), Inches(4.95),
               size=11.5, color=WHITE, dot_color=accent)

add_rect(s, 0, H - Inches(0.35), W, Inches(0.35), DARKER)
add_notes(s, """• 5-model ML ensemble: RF, Extra Trees, HGB, XGBoost, LightGBM then Ridge meta-learner
• OOF (out-of-fold) stacking prevents data leakage
• GeoPandas R-tree sindex makes spatial lookups O(log n)
• LRU cache means repeat queries skip the network entirely""")

# ─────────────────────────────────────────────
# SLIDE 4 — Price Modeling Architecture
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, LIGHT)
add_rect(s, 0, 0, W, Inches(1.2), NAVY)
add_rect(s, 0, Inches(1.2), Inches(0.07), H - Inches(1.2), YELLOW)

add_text(s, "Price Modeling Architecture",
         Inches(0.4), Inches(0.18), Inches(11), Inches(0.82),
         size=34, bold=True, color=WHITE)
add_text(s, "How we blend multiple signals into a single accurate estimate",
         Inches(0.4), Inches(0.84), Inches(11), Inches(0.35),
         size=13, color=YELLOW, italic=True)

fbox = s.shapes.add_shape(1, Inches(0.5), Inches(1.42), Inches(12.3), Inches(0.82))
fbox.line.fill.background()
fbox.fill.solid()
fbox.fill.fore_color.rgb = NAVY
add_text(s, "Final Price  =  p_base × 0.30  +  p_knn × 0.20  +  tract_anchor × 0.50",
         Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.65),
         size=19, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

steps = [
    ("p_base  ·  30%", BLUE, [
        "5-model OOF stacking ensemble",
        "RF + Extra Trees + HGB + XGBoost + LightGBM",
        "Ridge meta-learner on OOF predictions",
        "~30 engineered features per census tract",
        "Log-space prediction → exponentiated back",
        "Dampened to ±8% around market_base cap",
    ]),
    ("p_knn  ·  20%", VIOLET, [
        "7 nearest census tracts by lat/lon",
        "Median home value from ACS data",
        "Sklearn NearestNeighbors on centroids",
        "Captures hyper-local neighborhood pricing",
        "Bug fix Apr 2026: was never in the blend",
        "Now properly in the 3-way formula",
    ]),
    ("tract_anchor  ·  50%", GREEN, [
        "Direct census tract median home value",
        "Most weight — anchors to ground truth",
        "County appreciation factor (0.97–1.12×)",
        "Represents 2022 → 2024 market delta",
        "22 CA counties explicitly mapped",
        "Price-tiered fallback for unmapped counties",
    ]),
]

for i, (title, accent, items) in enumerate(steps):
    cx = Inches(0.35 + i * 4.3)
    card = s.shapes.add_shape(1, cx, Inches(2.42), Inches(4.1), Inches(4.72))
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    add_rect(s, cx, Inches(2.42), Inches(4.1), Inches(0.06), accent)
    add_text(s, title, cx + Inches(0.15), Inches(2.52), Inches(3.8), Inches(0.42),
             size=14, bold=True, color=accent)
    bullet_box(s, items, cx + Inches(0.12), Inches(2.98), Inches(3.85), Inches(4.0),
               size=11.5, color=NAVY, dot_color=accent)

add_rect(s, 0, H - Inches(0.4), W, Inches(0.4), NAVY)
add_notes(s, """• tract_anchor at 50% — raw census ground truth dominates
• p_knn fills spatial gaps between tracts
• p_base adds cross-CA ML pattern recognition
• When RentCast is available it overrides at 70%/30% on top of this blend""")

# ─────────────────────────────────────────────
# SLIDE 5 — Accuracy Bug Fixes
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.08), YELLOW)

add_text(s, "Accuracy Bug Fixes",
         Inches(0.5), Inches(0.22), Inches(10), Inches(0.82),
         size=36, bold=True, color=WHITE)
add_text(s, "Four critical fixes — silent failures that degraded every estimate",
         Inches(0.5), Inches(0.88), Inches(11), Inches(0.4),
         size=14, color=YELLOW, italic=True)

fixes = [
    ("KNN Self-Exclusion Bug", GREEN,
     "Query used idxs[0][1:] — a training-time self-exclusion pattern that silently "
     "skipped the closest census tract comp. Changed to idxs[0]. Nearest neighbor now "
     "correctly included."),
    ("KNN Result Never Used in Blend", YELLOW,
     "p_knn was computed but the formula only used p_base and tract_anchor (45/55 split). "
     "Changed to 3-way blend (30/20/50). Neighborhood comps now actually influence price."),
    ("Contribution Analysis County Factor", BLUE,
     "_price_with() re-computed market_base using a fallback formula, ignoring the "
     "resolved county lookup. Fixed to reuse pre-computed market_base — driver bars "
     "now accurate for all 22 mapped counties."),
    ("Coast Distance Cross-Validation", VIOLET,
     "Added KD-tree geometric check. If the census tract value and KD-tree disagree by "
     ">50%, the KD-tree overrides — eliminating stale tract data and shapely polygon "
     "interior measurement errors."),
]

for i, (title, accent, desc) in enumerate(fixes):
    row = i // 2
    col = i % 2
    cx = Inches(0.35 + col * 6.5)
    cy = Inches(1.48 + row * 2.78)
    card = s.shapes.add_shape(1, cx, cy, Inches(6.2), Inches(2.6))
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = DARKBG
    add_rect(s, cx, cy, Inches(6.2), Inches(0.06), accent)
    add_text(s, title, cx + Inches(0.15), cy + Inches(0.1), Inches(5.8), Inches(0.38),
             size=14, bold=True, color=accent)
    add_text(s, desc, cx + Inches(0.15), cy + Inches(0.54), Inches(5.9), Inches(1.9),
             size=11.5, color=WHITE, wrap=True)

add_rect(s, 0, H - Inches(0.35), W, Inches(0.35), DARKER)
add_notes(s, """• KNN fix: closest comp was being skipped every single query
• Blend fix: p_knn was computed but thrown away before the final number
• Contribution fix: driver bars were using wrong county factor
• Coast KD-tree: independent geometry check catches stale tract data""")

# ─────────────────────────────────────────────
# SLIDE 6 — 20 UI/UX Improvements
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, LIGHT)
add_rect(s, 0, 0, W, Inches(1.2), NAVY)
add_rect(s, 0, Inches(1.2), Inches(0.07), H - Inches(1.2), YELLOW)

add_text(s, "20 UI/UX Improvements",
         Inches(0.4), Inches(0.18), Inches(10), Inches(0.82),
         size=34, bold=True, color=WHITE)
add_text(s, "Quality-of-life upgrades shipped in a single sprint",
         Inches(0.4), Inches(0.84), Inches(10), Inches(0.35),
         size=13, color=YELLOW, italic=True)

bullet_box(s, [
    "Animated loading bar during API calls",
    "Click-to-copy price on valuation hero",
    "Browser tab title updates on each search",
    "/ keyboard shortcut focuses search bar",
    "Search history dropdown (last 6 searches)",
    "Chart.js dark mode palettes for all charts",
    "My Location button on map (Geolocation API)",
    "Tracker stats row: count, avg, highest value",
    "Housing cost warning if >30% of income",
    "Savings warning if monthly surplus < $0",
], Inches(0.5), Inches(1.48), Inches(6.1), Inches(5.7),
   size=13.5, color=NAVY, dot_color=YELLOW)

bullet_box(s, [
    "Transport budget field in LifeBudget tab",
    "Clickable market stats ticker on Home tab",
    "Analysis subtitle shows current address",
    "Print Report button on valuation result",
    "Input validation with inline error messages",
    "Tracker total animates on value change",
    "Empty state shows 3 example address pills",
    "Clear button resets the Live Map search",
    "Confidence range tooltip with explanation",
    "Dark mode active tab highlight fix",
], Inches(6.9), Inches(1.48), Inches(6.0), Inches(5.7),
   size=13.5, color=NAVY, dot_color=BLUE)

add_rect(s, 0, H - Inches(0.4), W, Inches(0.4), NAVY)
add_notes(s, """• 20 improvements in one sprint — focused on discoverability + feedback
• Loading bar + toasts give real-time status during slow API calls
• Tracker persistence means users keep their pipeline across sessions
• Dark mode was the most visible change""")

# ─────────────────────────────────────────────
# SLIDE 7 — Major Feature Additions
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.08), YELLOW)

add_text(s, "Major Feature Additions",
         Inches(0.5), Inches(0.22), Inches(10), Inches(0.82),
         size=36, bold=True, color=WHITE)
add_text(s, "Six substantial new capabilities added since April 17th",
         Inches(0.5), Inches(0.88), Inches(11), Inches(0.4),
         size=14, color=YELLOW, italic=True)

features = [
    ("Dark Mode", BLUE, [
        "Full-site Tailwind dark via html.dark class",
        "Persisted to localStorage",
        "Fixed chart palettes and text contrast",
    ]),
    ("Property Tracker", GREEN, [
        "Add/remove properties from any valuation",
        "Survives page refresh via localStorage",
        "Stats row: pipeline value, avg, highest",
    ]),
    ("Mortgage Breakdown", YELLOW, [
        "P+I, property tax (1.1%/yr), insurance",
        "PMI calculated if down payment < 20%",
        "Break-even vs. rent month shown",
    ]),
    ("Smart Search Routing", VIOLET, [
        "Regex /^\\d+\\s/ detects street addresses",
        "Address → Pro Analysis; text → Live Map",
        "Header + in-panel search use same logic",
    ]),
    ("In-Panel Address Search", ORANGE, [
        "Dedicated search bar in Pro Analysis panel",
        "No need to find the tiny header search",
        "Syncs with header search on every query",
    ]),
    ("RentCast Rent Estimate", GREEN, [
        "Calls /v1/avm/rent/long-term live",
        "Shows est. monthly rent + gross yield %",
        "Auto-fills rent comp for buy/rent chart",
    ]),
]

for i, (title, accent, bullets) in enumerate(features):
    row = i // 3
    col = i % 3
    cx = Inches(0.3 + col * 4.35)
    cy = Inches(1.42 + row * 2.88)
    card = s.shapes.add_shape(1, cx, cy, Inches(4.1), Inches(2.7))
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = DARKBG
    add_rect(s, cx, cy, Inches(4.1), Inches(0.05), accent)
    add_text(s, title, cx + Inches(0.12), cy + Inches(0.1), Inches(3.82), Inches(0.4),
             size=14, bold=True, color=accent)
    bullet_box(s, bullets, cx + Inches(0.1), cy + Inches(0.56), Inches(3.9), Inches(2.0),
               size=11.5, color=WHITE, dot_color=accent)

add_rect(s, 0, H - Inches(0.35), W, Inches(0.35), DARKER)
add_notes(s, """• Dark mode — persists across sessions, all components covered
• Tracker persistence — critical for real estate deal workflow
• Mortgage breakdown gives true all-in monthly cost, not just P&I
• Rent estimate auto-populates the buy vs. rent chart with real market data""")

# ─────────────────────────────────────────────
# SLIDE 8 — External Data Sources & APIs
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, LIGHT)
add_rect(s, 0, 0, W, Inches(1.2), NAVY)
add_rect(s, 0, Inches(1.2), Inches(0.07), H - Inches(1.2), YELLOW)

add_text(s, "Data Sources & External APIs",
         Inches(0.4), Inches(0.18), Inches(11), Inches(0.82),
         size=34, bold=True, color=WHITE)
add_text(s, "Open data + proprietary APIs combined for maximum accuracy",
         Inches(0.4), Inches(0.84), Inches(11), Inches(0.35),
         size=13, color=YELLOW, italic=True)

sources = [
    ("ACS 2021-22 Census", YELLOW, [
        "~8,000 California census tracts",
        "Median home value, income, age, population",
        "Housing units, population density",
        "Foundation of all ML training features",
    ]),
    ("TIGER Shapefiles", BLUE, [
        "Tract polygon boundaries for spatial joins",
        "R-tree sindex for O(log n) point-in-polygon",
        "CA coastline geometry for distance calc",
        "City centroids GeoDataFrame",
    ]),
    ("Nominatim / OpenStreetMap", GREEN, [
        "Free geocoder — address → lat/lon",
        "3-tier fallback: CA bounded → US → global",
        "LRU-cached per address string",
        "3-attempt exponential backoff retry",
    ]),
    ("RentCast AVM APIs", VIOLET, [
        "/v1/avm/value — comparable sales estimate",
        "/v1/avm/rent/long-term — rent estimate",
        "Blended at 70% when available",
        "LRU-cached per address + specs",
    ]),
    ("Google Street View Static", ORANGE, [
        "640×420 exterior photo per property",
        "FOV 90°, pitch 5° for street-level framing",
        "Outdoor source filter avoids indoor shots",
        "Shown in Pro Analysis left panel",
    ]),
    ("CA Coastline KD-Tree", RED, [
        "Independent coast distance verification",
        "Built from CA shoreline vertex coordinates",
        "EPSG:3310 projected — meter-level accuracy",
        "Overrides tract value if discrepancy >50%",
    ]),
]

for i, (title, accent, bullets) in enumerate(sources):
    row = i // 3
    col = i % 3
    cx = Inches(0.3 + col * 4.35)
    cy = Inches(1.42 + row * 2.82)
    card = s.shapes.add_shape(1, cx, cy, Inches(4.1), Inches(2.62))
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    add_rect(s, cx, cy, Inches(4.1), Inches(0.05), accent)
    add_text(s, title, cx + Inches(0.12), cy + Inches(0.1), Inches(3.85), Inches(0.4),
             size=13, bold=True, color=accent)
    bullet_box(s, bullets, cx + Inches(0.1), cy + Inches(0.54), Inches(3.9), Inches(1.98),
               size=11.5, color=NAVY, dot_color=accent)

add_rect(s, 0, H - Inches(0.4), W, Inches(0.4), NAVY)
add_notes(s, """• Census data is the backbone — 8k tracts across all of CA
• RentCast gives real MLS signal — overrides our model at 70% weight
• Nominatim is free but rate-limited to 1 req/sec — LRU cache is essential
• KD-tree is the safety net for coast distance edge cases""")

# ─────────────────────────────────────────────
# SLIDE 9 — ML Feature Engineering
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.08), YELLOW)

add_text(s, "ML Feature Engineering",
         Inches(0.5), Inches(0.22), Inches(10), Inches(0.82),
         size=36, bold=True, color=WHITE)
add_text(s, "~30 features per census tract feeding the 5-model stacking ensemble",
         Inches(0.5), Inches(0.88), Inches(11), Inches(0.4),
         size=14, color=YELLOW, italic=True)

feat_cols = [
    ("Raw Census Features", BLUE, [
        "median_income",
        "median_house_age",
        "population",
        "housing_units",
        "pop_density",
        "dist_coast_km",
        "dist_city_km",
        "lat, lon (tract centroid)",
    ]),
    ("Engineered Interactions", YELLOW, [
        "income_per_age = income / (age+1)",
        "coast_income = income / (coast+1)",
        "city_density = density / (dist_city+1)",
        "income_density = income x density / 1e8",
        "occupancy_pressure = pop / units",
        "relative_density = density / CA_avg",
        "log_income, log_density",
        "log_coast, log_city",
    ]),
    ("Geospatial & Polynomial", GREEN, [
        "lat_sq, lon_sq, lat_lon",
        "dist_sv — Silicon Valley distance",
        "dist_bh — Beverly Hills distance",
        "dist_lj — La Jolla distance",
        "dist_marin — Marin County distance",
        "coast_sq, age_sq (curvature)",
        "coast_city_ratio = dist_city/(coast+1)",
        "income_sq = income^2 / 1e9",
    ]),
]

for i, (title, accent, items) in enumerate(feat_cols):
    cx = Inches(0.35 + i * 4.3)
    card = s.shapes.add_shape(1, cx, Inches(1.42), Inches(4.1), Inches(5.72))
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = DARKBG
    add_rect(s, cx, Inches(1.42), Inches(4.1), Inches(0.07), accent)
    add_text(s, title, cx + Inches(0.15), Inches(1.54), Inches(3.8), Inches(0.44),
             size=14, bold=True, color=accent)
    bullet_box(s, items, cx + Inches(0.12), Inches(2.02), Inches(3.85), Inches(5.0),
               size=12, color=WHITE, dot_color=accent)

add_rect(s, 0, H - Inches(0.35), W, Inches(0.35), DARKER)
add_notes(s, """• Raw features alone miss non-linear value drivers
• coast_income: high-income coastal areas cost disproportionately more
• 4 luxury hotspot distances: Silicon Valley, Beverly Hills, La Jolla, Marin
• Log transforms match how income and density are actually distributed""")

# ─────────────────────────────────────────────
# SLIDE 10 — Results & Impact
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, LIGHT)
add_rect(s, 0, 0, W, Inches(1.2), NAVY)
add_rect(s, 0, Inches(1.2), Inches(0.07), H - Inches(1.2), YELLOW)

add_text(s, "Results & Impact",
         Inches(0.4), Inches(0.18), Inches(10), Inches(0.82),
         size=34, bold=True, color=WHITE)
add_text(s, "Summary of what was built and where we’re headed",
         Inches(0.4), Inches(0.84), Inches(10), Inches(0.35),
         size=13, color=YELLOW, italic=True)

stats = [
    ("4", "Accuracy Bug Fixes", YELLOW),
    ("20+", "UI/UX Improvements", BLUE),
    ("6", "Major New Features", GREEN),
    ("3-Way", "Price Blend Formula", VIOLET),
]

for i, (big, small, accent) in enumerate(stats):
    cx = Inches(0.4 + i * 3.2)
    card = s.shapes.add_shape(1, cx, Inches(1.48), Inches(3.0), Inches(1.75))
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY
    add_rect(s, cx, Inches(1.48), Inches(3.0), Inches(0.06), accent)
    add_text(s, big, cx + Inches(0.1), Inches(1.58), Inches(2.8), Inches(0.72),
             size=30, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(s, small, cx + Inches(0.1), Inches(2.2), Inches(2.8), Inches(0.5),
             size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Next Steps & Roadmap",
         Inches(0.5), Inches(3.45), Inches(6.2), Inches(0.45),
         size=16, bold=True, color=NAVY)
bullet_box(s, [
    "Expand RentCast rent AVM to auto-populate all searches",
    "Zillow / Redfin data scraping for additional price signal",
    "User accounts + cloud-synced property tracker",
    "Mobile-responsive layout for field use",
    "Export valuations to PDF / Excel for client presentations",
], Inches(0.5), Inches(3.95), Inches(6.0), Inches(3.1),
   size=13, color=NAVY, dot_color=YELLOW)

add_text(s, "Key Lesson",
         Inches(7.2), Inches(3.45), Inches(5.7), Inches(0.45),
         size=16, bold=True, color=NAVY)
add_text(s,
         "The biggest accuracy wins came from fixing silent bugs that caused "
         "real data to be ignored — not from adding new models. "
         "Correctness first, complexity second.",
         Inches(7.2), Inches(3.95), Inches(5.7), Inches(1.85),
         size=13, color=SLATE, italic=True, wrap=True)

add_rect(s, 0, H - Inches(0.4), W, Inches(0.4), NAVY)
add_notes(s, """• All 4 bugs were silent — code ran fine but used wrong data
• 20 UI improvements shipped in one batch = high leverage
• Next priority: mobile layout and cloud sync
• Lesson: always verify computed values are actually used in the final output""")

# ─────────────────────────────────────────────
# SLIDE 11 — Closing
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, W, Inches(0.08), YELLOW)
add_rect(s, 0, H - Inches(0.08), W, Inches(0.08), YELLOW)

add_text(s, "PriceNest",
         Inches(0.5), Inches(1.85), Inches(12), Inches(1.5),
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "AI-Powered  ·  California-Wide  ·  Built from Scratch",
         Inches(0.5), Inches(3.25), Inches(12), Inches(0.65),
         size=22, color=YELLOW, align=PP_ALIGN.CENTER, italic=True)
add_text(s, "Jake  ·  Saketh  ·  Aditiya",
         Inches(0.5), Inches(4.35), Inches(12), Inches(0.5),
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "April 2026",
         Inches(0.5), Inches(4.95), Inches(12), Inches(0.42),
         size=14, color=SLATE, align=PP_ALIGN.CENTER, italic=True)

add_notes(s, """• Thank the audience
• Open to Q&A
• Live demo available — can search any California address in real time""")

# ─────────────────────────────────────────────
# Save file
# ─────────────────────────────────────────────
out = "/Users/jakebeater/Downloads/comps_pricenest/PriceNest_April2026.pptx"
prs.save(out)
print(f"Saved: {out}")
