"""
Generate PriceNest summary presentation as a .pptx file.
Run: python3 generate_slides.py
Then upload the output file to Google Drive — it auto-converts to Google Slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
NAVY   = RGBColor(0x0B, 0x1C, 0x3E)
GOLD   = RGBColor(0xFD, 0xBB, 0x11)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xE8, 0xEE, 0xF8)
SLATE  = RGBColor(0x64, 0x74, 0x8B)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
RED    = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank layout

# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullets(slide, items, left, top, width, height,
                font_size=14, color=WHITE, indent=0.18, spacing_after=6):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(spacing_after)
        p.space_before = Pt(2)
        # Sub-bullet if starts with spaces
        is_sub = item.startswith("  ")
        bullet_char = "  ▸ " if is_sub else "• "
        run = p.add_run()
        run.text = bullet_char + item.strip()
        run.font.size = Pt(font_size - 1 if is_sub else font_size)
        run.font.color.rgb = LIGHT if is_sub else color
        run.font.bold = False
    return txBox

def navy_slide(title_text, subtitle_text=None):
    slide = prs.slides.add_slide(BLANK)
    # Full navy background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
    # Gold accent bar on left
    add_rect(slide, 0, 0, 0.25, 7.5, fill_rgb=GOLD)
    # Title
    add_text(slide, title_text, 0.55, 0.28, 12.0, 0.85,
             font_size=36, bold=True, color=WHITE)
    # Gold underline
    add_rect(slide, 0.55, 1.12, 2.5, 0.045, fill_rgb=GOLD)
    if subtitle_text:
        add_text(slide, subtitle_text, 0.55, 1.22, 11.5, 0.55,
                 font_size=16, color=LIGHT, italic=True)
    return slide

def two_col_slide(title_text, subtitle_text=None):
    slide = navy_slide(title_text, subtitle_text)
    return slide

# ── SLIDE 1 — Cover ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=GOLD)
add_rect(slide, 0.35, 3.18, 12.98, 0.065, fill_rgb=GOLD)

add_text(slide, "PriceNest", 0.65, 0.9, 12.0, 1.3,
         font_size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "AI-Powered California Home Valuation", 0.65, 2.25, 11.5, 0.7,
         font_size=24, color=GOLD, italic=True)
add_text(slide, "14-Day Development Summary", 0.65, 3.38, 10.0, 0.55,
         font_size=18, color=LIGHT)
add_text(slide, "By Jake, Saketh & Adih", 0.65, 6.5, 12.0, 0.65,
         font_size=15, color=SLATE, italic=True)
add_text(slide, "April 2025", 10.5, 6.5, 2.5, 0.65,
         font_size=15, color=SLATE, align=PP_ALIGN.RIGHT)

# ── SLIDE 2 — What is PriceNest? ─────────────────────────────────────────────
slide = navy_slide("What is PriceNest?", "The product we built")
add_bullets(slide, [
    "A web application that estimates California home values using AI and public data",
    "Users enter any CA address + property specs (sqft, beds, baths, condition)",
    "The model returns an estimated price, confidence range, and breakdown of price drivers",
    "",
    "Key capabilities:",
    "  Live Map — interactive hex heatmap of home values across all of California",
    "  Pro Analysis — buy vs. rent cost comparison with multi-year chart",
    "  Price Driver Breakdown — shows which factors are pushing the price up or down",
    "  Market Stats Panel — real-time CA median, coastal premium, and tier benchmarks",
], left=0.55, top=1.65, width=12.3, height=5.4, font_size=15)

# ── SLIDE 3 — Tech Stack & Data Sources ──────────────────────────────────────
slide = navy_slide("Tech Stack & Data Sources", "Resources powering PriceNest")

add_text(slide, "Backend", 0.55, 1.65, 5.8, 0.4, font_size=16, bold=True, color=GOLD)
add_bullets(slide, [
    "Python / Flask — web server and REST API",
    "scikit-learn — machine learning models",
    "GeoPandas + Shapely — spatial joins and geometry",
    "scipy cKDTree — fast nearest-neighbor coast distance",
    "geopy / Nominatim — address geocoding",
    "Flask-Limiter — API rate limiting",
], left=0.55, top=2.1, width=5.9, height=4.5, font_size=13)

add_text(slide, "Frontend & Data", 7.0, 1.65, 5.8, 0.4, font_size=16, bold=True, color=GOLD)
add_bullets(slide, [
    "Tailwind CSS — utility-first responsive styling",
    "Leaflet.js — interactive map and hex heatmap",
    "Chart.js — buy vs. rent comparison chart",
    "US Census ACS — income, density, housing data",
    "GSHHS Shoreline — high-resolution CA coastline",
    "Nominatim (OpenStreetMap) — geocoding service",
], left=7.0, top=2.1, width=5.9, height=4.5, font_size=13)

add_rect(slide, 6.7, 1.5, 0.06, 5.7, fill_rgb=GOLD)

# ── SLIDE 4 — ML Model Improvements ──────────────────────────────────────────
slide = navy_slide("ML Model Improvements", "From a single model to a full ensemble stack")
add_bullets(slide, [
    "Started with a single RandomForestRegressor as the base model",
    "",
    "Upgraded to a VotingRegressor ensemble (3 models averaged together):",
    "  RandomForestRegressor (300 trees, depth 15) — captures non-linear interactions",
    "  HistGradientBoostingRegressor (400 iters) — handles mixed feature scales",
    "  ExtraTreesRegressor (200 trees) — adds variance reduction and stability",
    "",
    "Added KNN stacking layer — finds 5 similar tracts and price-adjusts by square footage",
    "Ridge meta-model (α tuned from 1000 → 8) blends base model + KNN comps",
    "50% tract census anchor added — uses the neighborhood's own recorded median as a market floor",
    "Income-scaled appreciation factor (1.30×–1.65×) corrects for ACS 2020 → 2025 drift",
    "Trimmed top/bottom 1% of training prices to reduce extreme-outlier skew",
], left=0.55, top=1.65, width=12.3, height=5.4, font_size=14)

# ── SLIDE 5 — New Features: Interaction Variables ─────────────────────────────
slide = navy_slide("New Model Features", "Interaction variables that improved accuracy")

add_text(slide, "12 total features used in training:", 0.55, 1.65, 12.0, 0.4,
         font_size=15, color=LIGHT)

add_text(slide, "Core Census Features", 0.55, 2.15, 5.8, 0.4, font_size=15, bold=True, color=GOLD)
add_bullets(slide, [
    "median_income — household income in the tract",
    "median_house_age — average age of housing stock",
    "population & housing_units — tract demographics",
    "pop_density — people per square km",
    "dist_coast_km — miles to Pacific coastline",
    "dist_city_km — miles to nearest major city",
], left=0.55, top=2.6, width=5.9, height=4.0, font_size=13)

add_text(slide, "Engineered Interaction Features", 7.0, 2.15, 5.8, 0.4, font_size=15, bold=True, color=GOLD)
add_bullets(slide, [
    "income_per_age — wealth relative to housing vintage",
    "persons_per_unit — occupancy pressure signal",
    "relative_density — density vs. CA state average",
    "coast_income — 'Malibu effect': coastal wealth premium",
    "city_density — urban density pressure index",
], left=7.0, top=2.6, width=5.9, height=4.0, font_size=13)

add_rect(slide, 6.7, 2.0, 0.06, 5.0, fill_rgb=GOLD)

# ── SLIDE 6 — Ocean Distance Fix ─────────────────────────────────────────────
slide = navy_slide("Ocean Distance Feature Fix", "The most critical accuracy bug we resolved")
add_bullets(slide, [
    "Problem: coast distance was returning 0 for every inland property",
    "  Root cause: GSHHS files contain land-area polygons — a point inside the polygon has distance = 0",
    "  This made ALL predictions ignore the coast penalty entirely",
    "",
    "Fix 1 — Extract exterior boundary rings:",
    "  Converted filled land polygons → exterior line strings (the actual shoreline)",
    "  Shapely .exterior gives the true coastline edge for distance calculations",
    "",
    "Fix 2 — Independent KD-tree double-check:",
    "  Built a scipy cKDTree from all coastal vertex coordinates (EPSG:3310 meters)",
    "  If the primary distance and KD-tree distance disagree by >50%, the KD-tree wins",
    "  This completely eliminates silent failures from stale or bad tract data",
    "",
    "Result: Malibu now correctly shows near-zero coast distance; Sacramento ~120 km",
], left=0.55, top=1.65, width=12.3, height=5.4, font_size=14)

# ── SLIDE 7 — Geocoding Fix ───────────────────────────────────────────────────
slide = navy_slide("California-Only Geocoding", "Ensuring searches always resolve within CA")
add_bullets(slide, [
    "Problem: searching '4509 Camden Court' returned a Chicago suburb, not Granite Bay, CA",
    "  The app would then predict a negative or nonsensical price outside the model's training range",
    "",
    "3-tier California-preferring geocoder (both backend and frontend):",
    "",
    "  Tier 1 — Hard-bounded CA viewbox (bounded=1)",
    "    Nominatim only returns results within California's geographic bounding box",
    "",
    "  Tier 2 — Soft CA preference (bounded=0)",
    "    CA viewbox used as a ranking hint; first CA result accepted if one exists",
    "",
    "  Tier 3 — Auto-append ', California'",
    "    If no CA result found, the query is retried with ', California' appended",
    "",
    "CA bounds check in /predict endpoint rejects any non-CA coordinates outright",
    "Map tab search bar (JavaScript) mirrors the same 3-tier logic independently",
], left=0.55, top=1.65, width=12.3, height=5.4, font_size=13.5)

# ── SLIDE 8 — UI/UX Overhaul ─────────────────────────────────────────────────
slide = navy_slide("UI / UX Improvements", "From functional to polished")

add_text(slide, "Design System", 0.55, 1.65, 5.8, 0.4, font_size=15, bold=True, color=GOLD)
add_bullets(slide, [
    "Brand palette: Navy #0B1C3E + Gold #FDBB11",
    "Glass morphism cards with backdrop blur",
    "Plus Jakarta Sans — clean, modern typeface",
    "Tailwind CSS utility classes throughout",
    "fadeIn tab transitions (0.22s ease)",
    "Responsive layout for all screen sizes",
], left=0.55, top=2.1, width=5.9, height=4.5, font_size=13)

add_text(slide, "Interaction & Feedback", 7.0, 1.65, 5.8, 0.4, font_size=15, bold=True, color=GOLD)
add_bullets(slide, [
    "Skeleton loading cards before prediction loads",
    "Animated counter stats on homepage",
    "Price tier badges (Budget / Mid / Premium / Luxury)",
    "Toast notifications for errors and status",
    "Neighborhood rank bar (percentile within CA)",
    "Contribution bars sorted by dollar impact",
], left=7.0, top=2.1, width=5.9, height=4.5, font_size=13)

add_rect(slide, 6.7, 1.5, 0.06, 5.7, fill_rgb=GOLD)

# ── SLIDE 9 — New Tabs & Panels ───────────────────────────────────────────────
slide = navy_slide("New Tabs & Panels", "Features added during the 14-day sprint")
add_bullets(slide, [
    "Live Map Tab",
    "  Hex heatmap overlaid on Leaflet map — shows home values across all of California",
    "  Click any location to get an AI valuation card for that address",
    "  Search bar with 3-tier CA geocoding built into the map",
    "",
    "Pro Analysis Tab",
    "  Side-by-side monthly cost: buy (mortgage + taxes) vs. rent estimate",
    "  Multi-year cumulative cost chart (Chart.js) showing break-even point",
    "  Property spec controls: sqft, beds, baths, condition",
    "",
    "Price Driver Breakdown",
    "  Counterfactual analysis: each feature swapped to CA median to measure its impact",
    "  Top 5 drivers shown as sorted contribution bars with +/- dollar amounts",
    "",
    "Live CA Market Stats Panel",
    "  Statewide median, coastal premium %, 25th/75th/90th percentile benchmarks",
    "  Animated counters load at page open via /market-stats endpoint",
], left=0.55, top=1.65, width=12.3, height=5.4, font_size=13)

# ── SLIDE 10 — Accuracy Before vs After ──────────────────────────────────────
slide = navy_slide("Prediction Accuracy: Before vs. After", "Real-world example — Granite Bay, CA")

# Before box
add_rect(slide, 0.55, 1.65, 5.7, 5.0, fill_rgb=RGBColor(0x1E, 0x1E, 0x3E), line_rgb=RED, line_width_pt=1.5)
add_text(slide, "BEFORE", 0.75, 1.8, 5.0, 0.4, font_size=13, bold=True, color=RED)
add_bullets(slide, [
    "Single RandomForest, no stacking",
    "ACS 2020 data with no appreciation adjustment",
    "Ridge meta-model α=1000 pulling everything to CA average",
    "Coast distance bug returning 0 for inland homes",
    "Geocoding returning Illinois addresses",
    "",
    "Granite Bay estimated price:",
    "  ~$450,000",
    "  (actual 2025 median: $950K–$1.3M)",
], left=0.75, top=2.25, width=5.3, height=4.0, font_size=13, color=LIGHT)

# After box
add_rect(slide, 6.85, 1.65, 5.7, 5.0, fill_rgb=RGBColor(0x05, 0x2E, 0x1E), line_rgb=GREEN, line_width_pt=1.5)
add_text(slide, "AFTER", 7.05, 1.8, 5.0, 0.4, font_size=13, bold=True, color=GREEN)
add_bullets(slide, [
    "3-model VotingRegressor ensemble",
    "50% tract anchor + income-scaled 1.30–1.65× appreciation",
    "Ridge α=8 with full feature sensitivity",
    "GSHHS exterior rings + KD-tree coast verification",
    "3-tier CA geocoding on both backend and frontend",
    "",
    "Granite Bay estimated price:",
    "  ~$1,050,000",
    "  (accurate to current Granite Bay market)",
], left=7.05, top=2.25, width=5.3, height=4.0, font_size=13, color=LIGHT)

add_text(slide, "vs", 6.28, 3.8, 0.6, 0.6, font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── SLIDE 11 — Future Plans ───────────────────────────────────────────────────
slide = navy_slide("Future Plans", "Where PriceNest goes next")

add_text(slide, "Near-Term (0–3 months)", 0.55, 1.65, 5.8, 0.4, font_size=15, bold=True, color=GOLD)
add_bullets(slide, [
    "Integrate real-time MLS / Zillow API for live comparable sales",
    "Add school district ratings as a model feature",
    "Walk Score / Transit Score integration for urban properties",
    "Saved searches — user accounts to track multiple properties",
    "Mobile-optimized responsive redesign",
], left=0.55, top=2.1, width=5.9, height=4.5, font_size=13)

add_text(slide, "Long-Term Vision", 7.0, 1.65, 5.8, 0.4, font_size=15, bold=True, color=GOLD)
add_bullets(slide, [
    "Nationwide expansion beyond California",
    "Multi-family and commercial property support",
    "Historical price trend charts per neighborhood",
    "Investment ROI calculator (cap rate, cash-on-cash)",
    "Natural language search ('3bd near beach under $1M')",
    "LLM-powered property report generation",
], left=7.0, top=2.1, width=5.9, height=4.5, font_size=13)

add_rect(slide, 6.7, 1.5, 0.06, 5.7, fill_rgb=GOLD)

# ── SLIDE 12 — Closing ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
add_rect(slide, 0, 0, 13.33, 0.35, fill_rgb=GOLD)
add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=GOLD)

add_text(slide, "Thank You", 0, 2.0, 13.33, 1.5,
         font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "PriceNest — Built by Jake, Saketh & Adih", 0, 3.6, 13.33, 0.6,
         font_size=20, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, "Accurate California home valuations, powered by AI", 0, 4.3, 13.33, 0.5,
         font_size=15, color=LIGHT, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUTPUT = "PriceNest_Presentation.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print("Upload this file to Google Drive — it auto-converts to Google Slides.")
